"""Build the MoP '26 final-presentation deck as a real .pptx.

Generated rather than hand-built for one reason: the fabrication photographs arrive after the
slides do. Every image is resolved at build time, and anything missing is drawn as a labelled
dashed placeholder instead of failing, so the deck is presentable now and improves as photos
land. Re-run after dropping files into `examples/photos/`; the script prints what is still
missing.

    .venv\\Scripts\\python.exe make_deck.py

Numbers on the slides are read from the shipped fabrication report where possible rather than
retyped, so a re-solve cannot silently desynchronise the deck from the build.
"""
from __future__ import annotations

import json
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
PHOTOS = ROOT / "examples" / "photos"
BUILD = ROOT / "out_pearl3_30" / "v6"
OUT = ROOT / "MoP26_ShadowArt_final.pptx"

W, H = 13.333, 7.5                      # 16:9 in inches
MARGIN = 0.72
BODY_W = W - 2 * MARGIN

BG = RGBColor(0x12, 0x14, 0x18)         # near-black; the piece is about light and its absence
FG = RGBColor(0xF2, 0xF2, 0xF0)
MUTED = RGBColor(0x92, 0x96, 0x9C)
ACCENT = RGBColor(0xE8, 0xB5, 0x4A)     # the lamp colour from the 3D scene
GOOD = RGBColor(0x7C, 0xC5, 0x76)
BAD = RGBColor(0xD9, 0x6B, 0x5E)
PANEL = RGBColor(0x1C, 0x20, 0x26)      # card fill, one step up from the background

FONT = "Segoe UI"
MONO = "Consolas"

_missing: list[str] = []


# ---------------------------------------------------------------- primitives

def _solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, rgb=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _solid(s, rgb)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = 0.04
    s.text_frame.text = ""
    return s


def _runs(p, text, size, color, weight_color=None, font=None):
    """Split `**bold**` markup into runs so a key term can carry the accent colour."""
    for i, chunk in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.name = font or FONT
        strong = i % 2 == 1
        r.font.bold = strong
        r.font.color.rgb = (weight_color or ACCENT) if strong else color


def text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """`lines` is a list of dicts: text, size, color, space_before/after, line_spacing."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.space_before = Pt(ln.get("before", 0))
        p.space_after = Pt(ln.get("after", 6))
        p.line_spacing = ln.get("spacing", 1.05)
        _runs(p, ln["text"], ln.get("size", 14), ln.get("color", FG), ln.get("strong", ACCENT),
              ln.get("font"))
    return box


def slide(prs, title=None, kicker=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    if kicker:
        text(s, MARGIN, 0.42, BODY_W, 0.3,
             [{"text": kicker.upper(), "size": 11, "color": ACCENT}])
    if title:
        text(s, MARGIN, 0.72 if kicker else 0.6, BODY_W, 0.8,
             [{"text": title, "size": 30, "color": FG, "after": 0}])
        r = rect(s, MARGIN, 1.52 if kicker else 1.4, 1.1, 0.028, ACCENT, MSO_SHAPE.RECTANGLE)
        r.shadow.inherit = False
    return s


def notes(s, body):
    s.notes_slide.notes_text_frame.text = body.strip()


# ---------------------------------------------------------------- images

def _find(stem):
    for ext in (".jpg", ".jpeg", ".png", ".JPG", ".PNG"):
        p = PHOTOS / f"{stem}{ext}"
        if p.exists():
            return p
    return None


def _fit(path, x, y, w, h):
    """Letterbox an image inside a box, preserving aspect ratio."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    nw, nh = (h * ar, h) if w / h > ar else (w, w / ar)
    return x + (w - nw) / 2, y + (h - nh) / 2, nw, nh


def picture(s, path, x, y, w, h, caption=None, label=None):
    """Insert `path`, or a dashed placeholder naming the file that is still missing."""
    p = Path(path) if path and Path(path).is_absolute() else (ROOT / path if path else None)
    if p is None or not p.exists():
        ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        _solid(ph, PANEL)
        ph.adjustments[0] = 0.03
        ph.line.color.rgb = ACCENT
        ph.line.width = Pt(1.25)
        ph.line.dash_style = 4                      # msoLineDash
        ph.text_frame.text = ""
        name = label or (p.name if p else "image")
        _missing.append(name)
        psz = 9 if w < 3.4 else 11                  # the path must not wrap mid-filename
        text(s, x + 0.12, y + h / 2 - 0.45, w - 0.24, 0.9,
             [{"text": "DROP PHOTO HERE", "size": 13, "color": ACCENT, "after": 4},
              {"text": f"examples/photos/{name}", "size": psz, "color": MUTED, "after": 4},
              {"text": caption or "", "size": 10, "color": MUTED}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    fx, fy, fw, fh = _fit(p, x, y, w, h if not caption else h - 0.3)
    s.shapes.add_picture(str(p), Inches(fx), Inches(fy), Inches(fw), Inches(fh))
    if caption:
        text(s, x, y + h - 0.26, w, 0.26,
             [{"text": caption, "size": 10, "color": MUTED}], align=PP_ALIGN.CENTER)


def photo(s, stem, x, y, w, h, caption=None):
    picture(s, _find(stem), x, y, w, h, caption=caption, label=f"{stem}.jpg")


# ---------------------------------------------------------------- composites

def cards(s, y, items, height=1.55, top_size=26, gap=0.26):
    """Row of metric cards: (value, label) pairs."""
    n = len(items)
    w = (BODY_W - gap * (n - 1)) / n
    for i, (value, label) in enumerate(items):
        x = MARGIN + i * (w + gap)
        rect(s, x, y, w, height)
        text(s, x + 0.22, y + 0.22, w - 0.44, height - 0.44,
             [{"text": value, "size": top_size, "color": ACCENT, "after": 5},
              {"text": label, "size": 11.5, "color": MUTED, "spacing": 1.15}],
             anchor=MSO_ANCHOR.MIDDLE)


def verdict(s, x, y, w, tag, colour):
    b = rect(s, x, y, w, 0.28, PANEL, MSO_SHAPE.ROUNDED_RECTANGLE)
    b.line.color.rgb = colour
    b.line.width = Pt(0.75)
    text(s, x, y + 0.045, w, 0.24, [{"text": tag, "size": 9.5, "color": colour}],
         align=PP_ALIGN.CENTER)


def grid_table(s, x, y, w, rows, col_w, head=True, row_h=0.34):
    """Lightweight table: `rows` is a list of lists of strings."""
    for ri, row in enumerate(rows):
        yy = y + ri * row_h
        if head and ri == 0:
            rect(s, x, yy, w, row_h, PANEL, MSO_SHAPE.RECTANGLE)
        elif ri % 2 == 0:
            rect(s, x, yy, w, row_h, RGBColor(0x17, 0x1A, 0x1F), MSO_SHAPE.RECTANGLE)
        cx = x
        for ci, cell in enumerate(row):
            cw = col_w[ci]
            col = MUTED if (head and ri == 0) else FG
            size = 10.5 if (head and ri == 0) else 12
            text(s, cx + 0.14, yy + (row_h - 0.22) / 2, cw - 0.28, 0.24,
                 [{"text": cell, "size": size, "color": col}],
                 align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)
            cx += cw
    return y + len(rows) * row_h


def template_banner(s):
    b = rect(s, MARGIN, H - 0.92, BODY_W, 0.5, RGBColor(0x3A, 0x2A, 0x12))
    b.line.color.rgb = ACCENT
    b.line.width = Pt(1)
    text(s, MARGIN, H - 0.79, BODY_W, 0.3,
         [{"text": "TEMPLATE SLIDE — REPLACE BEFORE PRESENTING", "size": 12, "color": ACCENT}],
         align=PP_ALIGN.CENTER)


def divider(prs, title, sub):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    text(s, MARGIN, 2.95, BODY_W, 1.6,
         [{"text": title, "size": 38, "color": FG, "after": 16},
          {"text": sub, "size": 13.5, "color": MUTED}], align=PP_ALIGN.CENTER)
    r = rect(s, W / 2 - 0.55, 4.72, 1.1, 0.028, ACCENT, MSO_SHAPE.RECTANGLE)
    r.shadow.inherit = False
    return s


# ---------------------------------------------------------------- live numbers

def load_numbers():
    """Read headline metrics from the shipped report so the deck cannot drift from the build."""
    d = {"mean_iou": 0.846, "min_iou": 0.836, "ssim": 0.751, "rmse": 0.1405,
         "edge": 0.570, "stage1": 0.768, "stage2": 0.802, "slot": 3.46, "angle": 60.0}
    rep = BUILD / "fab_report.json"
    if not rep.exists():
        return d
    r = json.loads(rep.read_text())
    views = r["gate_fab_round_trip"]["views"]
    summ = r["gate_fab_round_trip"]["summary"]
    d["mean_iou"] = summ["mean_iou"]
    d["min_iou"] = summ["min_iou"]
    d["ssim"] = sum(v["ssim"] for v in views.values()) / len(views)
    d["rmse"] = sum(v["rmse"] for v in views.values()) / len(views)
    d["edge"] = sum(v["edge_fidelity"] for v in views.values()) / len(views)
    # The descent log carries BOTH numbers, and the gap between them is the slide's point:
    # the composite objective moves ~0.034 while mean IoU moves ~0.002, because IoU scores
    # silhouettes and the whole failure being fixed here is interior tone.
    d["descent"] = r["retone_descent"]
    d["score1"] = r["retone_descent"][0]["score"]
    d["score2"] = r["retone_descent"][-1]["score"]
    d["iou1"] = r["retone_descent"][0]["mean_iou"]
    d["iou2"] = r["retone_descent"][-1]["mean_iou"]
    d["slot"] = r["gate_weave"]["max_slot_width_mm"]
    d["angle"] = r["gate_weave"]["min_crossing_angle_deg"]
    return d


N = load_numbers()


# ================================================================= the deck

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)

    # ---- 1. title -------------------------------------------------------
    s = slide(prs)
    photo(s, "fab_04_object", 7.35, 0.85, 5.28, 5.8, None)
    text(s, MARGIN, 1.85, 6.3, 4,
         [{"text": "SHADOW ART", "size": 12, "color": ACCENT, "after": 14},
          {"text": "One object.", "size": 40, "color": FG, "after": 2},
          {"text": "Three portraits.", "size": 40, "color": FG, "after": 2},
          {"text": "Six sheets of acrylic.", "size": 40, "color": ACCENT, "after": 22},
          {"text": "A single laser-engraved assembly that casts the Girl with a Pearl "
                   "Earring from the front, the side and the back as it turns.",
           "size": 14, "color": MUTED, "spacing": 1.3, "after": 20},
          {"text": "Ofek Yankis  ·  Shadha  ·  Yahel Hershko", "size": 13, "color": FG,
           "after": 4},
          {"text": "MoP '26 Final Project  ·  Faculty of Architecture and Town Planning  "
                   "·  6 August 2026", "size": 11, "color": MUTED}])
    notes(s, """
[0:00-0:30] Open with the object in hand, not the slide.

"This is one object. It is six sheets of clear acrylic, engraved with a laser, slotted
together. On its own it looks like noise. Put a lamp behind it and it throws a portrait on
the wall. Turn it 120 degrees and you get the same person from the side. Turn it again and
you get the back of her head."

Then start the deck. Do not read the title slide aloud.
""")

    # ---- 2. team --------------------------------------------------------
    s = slide(prs, "The team", "Who we are")
    people = [
        ("Ofek Yankis", "BSc Computer Science, 3rd year",
         "Built the simulation environment and the differentiable renderer; GPU port and "
         "performance work; parameter sweeps and the fabrication export pipeline.",
         "Interests: [ADD — e.g. graphics, optimisation, GPU computing]"),
        ("Shadha", "BSc Computer Science, 4th year",
         "Designed the cross-talk-aware assignment algorithm and the two-stage solver; "
         "the damage functional and the re-toning descent.",
         "Interests: [ADD — e.g. algorithms, inverse problems, computational design]"),
        ("Yahel Hershko", "MSc Architecture",
         "Originated the concept; drove the physical and optical design; specified the "
         "material and joinery; fabricated and assembled the artefact.",
         "Interests: [ADD — e.g. digital fabrication, computational design, light]"),
    ]
    cw = (BODY_W - 0.52) / 3
    for i, (name, status, role, interests) in enumerate(people):
        x = MARGIN + i * (cw + 0.26)
        rect(s, x, 1.95, cw, 4.5)
        text(s, x + 0.28, 2.25, cw - 0.56, 3.9,
             [{"text": name, "size": 19, "color": FG, "after": 4},
              {"text": status, "size": 11.5, "color": ACCENT, "after": 16},
              {"text": role, "size": 12, "color": FG, "spacing": 1.28, "after": 16},
              {"text": interests, "size": 11, "color": MUTED, "spacing": 1.25}])
    text(s, MARGIN, 6.62, BODY_W, 0.4,
         [{"text": "Technion — Israel Institute of Technology  ·  "
                   "the project is a collaboration between Computer Science and Architecture, "
                   "which is also why it exists: the algorithm and the artefact were designed "
                   "against each other rather than in sequence.",
           "size": 11, "color": MUTED}])
    notes(s, """
[0:30-1:00] Thirty seconds. One line each.

The point worth landing: this is a CS/Architecture collaboration, and the fabrication
constraints were in the optimiser from the beginning rather than applied afterwards. That is
what the third bullet on the results slide is about.

TODO BEFORE PRESENTING: fill in the three "Interests" lines and Shadha's surname.
""")

    # ---- 3. the concept -------------------------------------------------
    s = slide(prs, "Three views of one painting, from one object", "Background — the idea")
    text(s, MARGIN, 2.0, 5.5, 4.2,
         [{"text": "We take Vermeer's **Girl with a Pearl Earring** and three views of "
                   "her — **front, profile, back** — and ask for a single physical object "
                   "whose shadow is all three, depending on how it is turned.",
           "size": 15, "color": FG, "spacing": 1.35, "after": 18},
          {"text": "The object stays on a turntable. The lamp never moves. Only the "
                   "rotation changes.", "size": 15, "color": FG, "spacing": 1.35,
           "after": 18},
          {"text": "That constraint is the whole problem. Three images have to be "
                   "encoded into one arrangement of matter, and every part of that matter "
                   "is visible to **all three** projections at once.",
           "size": 15, "color": FG, "spacing": 1.35}])
    picture(s, BUILD / "deliverables" / "pearl6_projections_preview.png",
            6.55, 1.9, 6.1, 4.6,
            caption="Left: the source at each of the three stops.  "
                    "Right: the shadow the six-sheet solution casts there.")
    notes(s, """
[1:00-1:40]

"Vermeer's Girl with a Pearl Earring. We want three views of her — front, profile, back —
out of one object on a turntable, with a fixed lamp.

The image on the right is not a photograph and it is not a mock-up. It is what our solution
computes the wall will look like, from the same geometry we then cut."

Do not explain the algorithm yet. Just establish the goal.
""")

    # ---- 4. how we got here ---------------------------------------------
    # The semester actually ran ceiling mobile -> Rhino/GH -> own simulator -> flat panels,
    # and the last step is what created the assignment problem the rest of the talk is about.
    s = slide(prs, "How we got here", "Background — the route")
    arc = [
        ("Hang it from the ceiling",
         "The first idea was a suspended cloud: fragments floating at their own depths, "
         "free in all three dimensions. Maximum freedom for the optimiser.",
         "Every fragment needs something to hang from, and the thread casts a shadow too. "
         "It was freedom the fabrication could not honour."),
        ("Rhino and Grasshopper",
         "We built it parametrically in Grasshopper, which is the standard tool for this "
         "kind of geometry in architecture.",
         "A dataflow evaluator, not an optimiser. We needed hundreds of solves, gradients "
         "and a GPU — and it could not show all three projections at once."),
        ("So we built our own simulator",
         "A differentiable renderer in Python, and an orbitable 3D scene in the browser — "
         "sheets, lamp, rays and all three walls in one view, with no CAD licence.",
         "The renderer became the physics model, the optimisation objective and the preview "
         "at the same time. Sweeps went from unrunnable to **21.6 s** per complete build."),
        ("Then the geometry changed completely",
         "The suspended cloud became **six flat laser-cut sheets**, slotted into each other: "
         "self-supporting, reproducible from files, nothing extra in the light path.",
         "But depth stopped being continuous. A fragment now has to pick a **discrete host "
         "sheet** — and that single change is what created the problem the rest of this "
         "talk is about."),
    ]
    cw = (BODY_W - 3 * 0.22) / 4
    for i, (head, what, learned) in enumerate(arc):
        x = MARGIN + i * (cw + 0.22)
        c = rect(s, x, 1.95, cw, 4.5)
        if i == 3:
            c.line.color.rgb = ACCENT
            c.line.width = Pt(1.25)
        text(s, x + 0.26, 2.18, cw - 0.52, 2.1,
             [{"text": f"0{i + 1}", "size": 11, "color": ACCENT, "after": 9},
              {"text": head, "size": 15, "color": FG, "spacing": 1.12, "after": 12},
              {"text": what, "size": 11.5, "color": FG, "spacing": 1.3}])
        text(s, x + 0.26, 4.62, cw - 0.52, 1.6,
             [{"text": "What it cost us", "size": 9.5, "color": ACCENT, "after": 6},
              {"text": learned, "size": 11, "color": MUTED, "spacing": 1.3}])
        if i < 3:
            text(s, x + cw + 0.01, 4.0, 0.2, 0.3,
                 [{"text": "›", "size": 15, "color": MUTED}], align=PP_ALIGN.CENTER)
    text(s, MARGIN, 6.62, BODY_W, 0.4,
         [{"text": "Two of the three big decisions this semester were made by discovering "
                   "what the previous one could not do.", "size": 12, "color": ACCENT}])
    notes(s, """
[1:40-2:40] Yahel opens this one — it is the architecture half of the story.

"We did not start here. We started wanting to hang fragments from the ceiling, and we built
it in Rhino and Grasshopper because that is what you use.

Two things pushed us off that. Grasshopper is a dataflow evaluator — it recomputes a graph.
It is not an optimiser, it has no GPU and no gradients, and we needed to run hundreds of
solves. So we wrote our own renderer and our own 3D viewer in the browser, and stopped
needing a CAD licence at all.

And once we could actually see the thing, the hanging design lost. It became six flat sheets
that slot into each other — no thread, no hardware in the beam, and the whole object comes
off a laser bed.

That last move is the one that matters for the algorithm, because depth stopped being a
continuous variable. A fragment now picks a discrete sheet. That is the problem we spent the
rest of the semester on."

Do not apologise for the dead ends. This is the "main challenges" item on the brief.
""")

    # ---- 5. prior art + contribution -------------------------------------
    s = slide(prs, "Where this sits, and what we added", "Background — prior art")
    entries = [
        ("Mitra & Pauly — Shadow Art", "SIGGRAPH Asia 2009", "NOT APPLICABLE", BAD,
         "Carve one connected volume from the visual hull. Ours is flat laser-cut stock, not "
         "a solid — but we kept their honesty: for many target triples no exact solution "
         "exists, so the goal is the best approximation, not a match."),
        ("Baran et al. — Layered Attenuators", "Eurographics 2012", "ADOPTED", GOOD,
         "The closest relative. Transmittances multiply, so optical densities add — tone "
         "levels should be spaced evenly in density. Worth **+0.034 mean IoU** to us. It then "
         "reversed when we cut to six sheets, which is the closing insight of this talk."),
        ("Bermano et al. — ShadowPix", "2012", "NOT APPLICABLE", BAD,
         "Height-field reliefs that self-shadow under grazing light. That is reflection; "
         "ours is transmission, so our stack composes multiplicatively instead."),
    ]
    for i, (title_, venue, tag, colour, why) in enumerate(entries):
        y = 1.95 + i * 1.56
        rect(s, MARGIN, y, 6.15, 1.48)
        text(s, MARGIN + 0.28, y + 0.2, 3.5, 0.5,
             [{"text": title_, "size": 13.5, "color": FG, "after": 2},
              {"text": venue, "size": 10, "color": MUTED}])
        verdict(s, MARGIN + 4.15, y + 0.24, 1.72, tag, colour)
        text(s, MARGIN + 0.28, y + 0.8, 5.6, 0.6,
             [{"text": why, "size": 10.5, "color": MUTED, "spacing": 1.25}])
    rect(s, 7.15, 1.95, 5.46, 4.6)
    items = [
        ("Rotation, not relighting",
         "The lamp is fixed and the object turns, so all three images share one set of "
         "matter with no ability to mask any of it."),
        ("Engraving density as the tone alphabet",
         "Four levels cut into one clear sheet. No dye, no lamination, no thickness "
         "variation — tone is a laser parameter."),
        ("Cross-talk priced, not suppressed",
         "Interference enters the objective as a **signed** term that may go negative, so a "
         "fragment is rewarded for genuinely serving another view."),
        ("Fabrication inside the loop",
         "Kerf, minimum feature and joint feasibility constrain the solver directly, and the "
         "gates re-verify against the exported polygons."),
    ]
    text(s, 7.45, 2.15, 4.9, 0.35, [{"text": "What is new here", "size": 14, "color": ACCENT}])
    yy = 2.66
    for head, body in items:
        text(s, 7.45, yy, 4.9, 0.9,
             [{"text": head, "size": 12.5, "color": FG, "after": 4},
              {"text": body, "size": 10.5, "color": MUTED, "spacing": 1.26}])
        yy += 0.97
    notes(s, """
[2:40-3:20] Forty seconds. Do not read the cards.

Say only this: "Two of these are not applicable and one is — Baran. From Baran we take that
the physics is multiplicative, and that because transmittances multiply, densities add, so
your tone levels should be evenly spaced in density. That was worth about three IoU points.
Hold that thought, because it reversed on us later."

Then point at the right-hand column and give one sentence: "what is new is that the lamp
never moves — the object turns — and that we price the interference between views instead of
trying to suppress it."
""")

    # ---- 6. why it is hard ----------------------------------------------
    s = slide(prs, "Why this is hard", "The problem")
    obstacles = [
        ("You cannot switch a sheet off",
         "Light does not care which picture a sheet was meant for. When the lamp is lit for "
         "the front view, **all six sheets** are in the beam — including the four intended "
         "for the other two views. Their shadows land on the wall too, stretched and smeared "
         "across the face. There is no masking operation available."),
        ("Shadows multiply, they do not add",
         "Two half-dark patches in line do not make black. Each passes half the light, so "
         "together they pass a quarter. Three sheets that each independently wanted mid-grey "
         "compose to **0.6³ = 0.22** — near black. Every sheet's contribution depends on what "
         "every other sheet does at that exact point."),
        ("No single sheet may be readable",
         "If one sheet carried a recognisable face, the piece would be a picture on plastic "
         "with decoration around it. The image has to exist **only in the combination**, so "
         "each portrait is deliberately shattered and scattered across sheets and depths."),
    ]
    cw = (BODY_W - 0.52) / 3
    for i, (head, body) in enumerate(obstacles):
        x = MARGIN + i * (cw + 0.26)
        rect(s, x, 1.98, cw, 4.4)
        text(s, x + 0.3, 2.28, cw - 0.6, 3.8,
             [{"text": f"0{i + 1}", "size": 13, "color": ACCENT, "after": 12},
              {"text": head, "size": 17, "color": FG, "spacing": 1.1, "after": 16},
              {"text": body, "size": 12.5, "color": MUTED, "spacing": 1.35}])
    notes(s, """
[3:20-4:20] Slow down here. If the audience gets these three, everything after is easy.

Obstacle 2 is the one to demonstrate with your hands — two flat palms, "half and half does
not make black, it makes three-quarters."

Land it: "so the naive design — one portrait per sheet — produces a black blob. We built it.
It does."
""")

    # ---- 7. the reframe -------------------------------------------------
    s = slide(prs, "The move: stop fighting the interference, and spend it", "The idea")
    text(s, MARGIN, 1.95, 6.2, 3.4,
         [{"text": "A fragment of engraving sitting in the beam casts a shadow for "
                   "**every** view, not just its own. Normally that is the error term.",
           "size": 15, "color": FG, "spacing": 1.35, "after": 16},
          {"text": "But if it is positioned so that its stray shadow lands where another "
                   "portrait **also wanted darkness**, that shadow is not damage. It is free "
                   "work. One piece of engraving, three pictures served.",
           "size": 15, "color": FG, "spacing": 1.35, "after": 16},
          {"text": "So the question the optimiser answers is not \"where do the front-view "
                   "pieces go\". It is: **where can a piece help all three views at once?**",
           "size": 15, "color": ACCENT, "spacing": 1.35}])
    picture(s, ROOT / "examples" / "pearl3" / "1_result.png", 7.15, 1.9, 5.5, 3.4)
    cards(s, 5.62, [("100%", "of fragments serve at least two of the three views"),
                    ("93%", "serve all three views"),
                    ("2.93 / 3", "average views served per fragment"),
                    ("−0.005", "measured cross-talk — negative, i.e. constructive")],
          height=1.35, top_size=24)
    notes(s, """
[4:20-5:05] This is the conceptual centre of the talk.

"Every piece of engraving casts a shadow for every view. That is normally the error term.
But if a piece lands where another portrait also wanted darkness, that stray shadow is free
work.

So we stopped asking 'where do the front-view pieces go' and started asking 'where can a
piece help all three views at once'."

Then the numbers: "93% of our fragments serve all three portraits. And the measured
cross-talk term is negative — the interference between the views makes the result better
than solving them in isolation."

That last number is the one to say slowly.
""")

    # ---- 8. pipeline ----------------------------------------------------
    s = slide(prs, "The pipeline", "The project — algorithmic approach")
    stages = [
        ("01", "Targets", "Three views become darkness fields on the wall"),
        ("02", "Optics", "Homographies, magnification, penumbra bound"),
        ("03", "Shatter", "Detail-adaptive Voronoi fragmentation"),
        ("04", "Assign", "Cross-talk-aware host selection — the core"),
        ("05", "Weave", "Half-lap joints, collision clearance"),
        ("06", "Re-tone", "Coordinate descent over four tones"),
        ("07", "Cut", "Min-feature, kerf, layered DXF/SVG"),
    ]
    cw = (BODY_W - 6 * 0.16) / 7
    for i, (num, head, body) in enumerate(stages):
        x = MARGIN + i * (cw + 0.16)
        highlight = num in ("04", "06")
        c = rect(s, x, 2.15, cw, 2.5, PANEL)
        if highlight:
            c.line.color.rgb = ACCENT
            c.line.width = Pt(1.25)
        text(s, x + 0.18, 2.4, cw - 0.36, 2.0,
             [{"text": num, "size": 11, "color": ACCENT if highlight else MUTED, "after": 9},
              {"text": head, "size": 14, "color": FG, "after": 9},
              {"text": body, "size": 10.5, "color": MUTED, "spacing": 1.25}])
        if i < 6:
            text(s, x + cw + 0.005, 3.2, 0.15, 0.3,
                 [{"text": "›", "size": 14, "color": MUTED}], align=PP_ALIGN.CENTER)
    text(s, MARGIN, 5.15, BODY_W, 1.4,
         [{"text": "One renderer serves as the physics simulator, the optimisation "
                   "objective and the preview.", "size": 14, "color": FG, "after": 8},
          {"text": "That identity is deliberate: it removes the entire class of bug where "
                   "the thing being optimised quietly diverges from the thing being "
                   "measured. It is differentiable, runs on the GPU, and a complete build "
                   "takes **21.6 seconds**.",
           "size": 12.5, "color": MUTED, "spacing": 1.3}])
    text(s, MARGIN, 6.55, BODY_W, 0.4,
         [{"text": "The two highlighted stages are where the research contribution sits.",
           "size": 11.5, "color": ACCENT}])
    notes(s, """
[5:05-5:35] Do not narrate all seven. Point at the chain, say "seven stages, targets in,
cut files out", then go straight to:

"Two of these are where the actual contribution is — assignment and re-toning. I want to
spend the next two and a half minutes on those."
""")

    # ---- 9. the algorithm (core) ----------------------------------------
    # The optics and the full damage functional now live in the backup section; this slide
    # has to survive being given ninety seconds and nothing else.
    s = slide(prs, "The algorithm: which sheet hosts each fragment",
              "The project — the contribution")
    text(s, MARGIN, 1.9, 6.55, 0.5,
         [{"text": "A fragment's position inside its own portrait is already fixed by the "
                   "optics. The one free variable left is **which sheet carries it** — and "
                   "that decides where its unavoidable stray shadow lands on the other two "
                   "views.", "size": 12.5, "color": FG, "spacing": 1.3}])
    rect(s, MARGIN, 2.85, 6.55, 3.28)
    code = [
        "shatter each view's target into shards",
        "order shards by the darkness they demand, densest first",
        "",
        "for each shard f:",
        "    for each sheet p that can physically reach f:",
        "        gain   \u2190 darkness f delivers to its OWN view",
        "        **damage \u2190 stray darkness it throws on the OTHER two**",
        "        score  \u2190 gain \u2212 \u03bb \u00b7 damage",
        "    host(f) \u2190 best-scoring sheet that still has room",
        "",
        "resolve slot collisions where sheets cross",
        "freeze the geometry; coordinate-descend the four tones",
    ]
    text(s, MARGIN + 0.3, 3.05, 6.0, 2.95,
         [{"text": ln or " ", "size": 10.5, "color": FG if ln.startswith(("for", "shatter",
                                                                          "order", "resolve",
                                                                          "freeze")) else MUTED,
           "font": MONO, "spacing": 1.18, "after": 2} for ln in code])
    text(s, MARGIN, 6.28, 6.55, 0.85,
         [{"text": "Before this, the host was a **random draw**.", "size": 12,
           "color": ACCENT, "after": 5},
          {"text": "Reconstruction quality swung **14–37% on the seed alone** — the single "
                   "largest source of variance in the project.",
           "size": 11.5, "color": MUTED, "spacing": 1.28}])
    rect(s, 7.42, 1.9, 5.19, 3.15)
    text(s, 7.72, 2.15, 4.6, 2.7,
         [{"text": "Damage is the expensive term — and we never render a wall to get it",
           "size": 13, "color": ACCENT, "spacing": 1.15, "after": 10},
          {"text": "Scoring a candidate naively means re-rendering another view's wall, for "
                   "every fragment and every sheet. Instead we compose the two homographies "
                   "once per (sheet, view) pair:",
           "size": 11.5, "color": MUTED, "spacing": 1.28, "after": 9},
          {"text": "G = H(sheet \u2192 other wall) \u00b7 H(own wall \u2192 sheet)",
           "size": 11, "color": ACCENT, "font": MONO, "after": 9},
          {"text": "That maps a point on the fragment's own wall **directly** to where its "
                   "material lands on another view's wall. Evaluating a host is then at most "
                   "**200 pixel transforms**.",
           "size": 11.5, "color": FG, "spacing": 1.28}])
    grid_table(s, 7.42, 5.25, 5.19,
               [["Cost of host selection", "complexity"],
                ["Naive — re-render per candidate", "F · P · V · N"],
                ["Ours — composed homography", "F · P · V · 200"],
                ["Speed-up on the dominant term", "≈ 1800×"]],
               [3.49, 1.70], row_h=0.32)
    text(s, 7.42, 6.62, 5.19, 0.5,
         [{"text": "The difference between choosing hosts by **optimisation** and choosing "
                   "them by coin flip.", "size": 11.5, "color": ACCENT, "spacing": 1.25}])
    notes(s, """
[5:35-7:05] THE slide. Ninety seconds, and do not rush it. Shadha takes this.

Four beats:
1. "The fragment's position in its own picture is already fixed by the optics. The only thing
   still free is which sheet carries it — and that decides where its stray shadow falls on
   the other two views."
2. Walk the pseudocode with a finger. Gain, damage, argmax. Say "the whole contribution is
   the one bold line."
3. "In our first version the host was a random draw. Quality swung 14 to 37 percent on the
   seed alone. It was random because scoring it properly means re-rendering a wall for every
   fragment, every sheet, every other view — far too expensive."
4. "So we never render. We compose the two homographies into one matrix that maps directly
   from one wall to another, and transform a couple of hundred of the fragment's own pixels.
   About eighteen hundred times cheaper — and that is the difference between optimising the
   choice and guessing it."

THERE IS A BACKUP SLIDE with the damage functional written out, its four limiting cases, and
why greedy. Go to it if anyone asks how damage is actually defined, or whether the greedy is
optimal. Do not pre-empt it here.
""")

    # ---- 11. re-toning --------------------------------------------------
    s = slide(prs, "The second stage, and the metric that hid the bug", "The project — the contribution")
    text(s, MARGIN, 1.9, 6.0, 4.6,
         [{"text": f"Our renders scored {N['iou1']:.3f} on IoU and still looked wrong.",
           "size": 15, "color": ACCENT, "after": 10},
          {"text": "Both were true. **IoU scores the outline of a shape**, so a face that is "
                   "correctly shaped and uniformly crushed to black scores well. The "
                   "silhouettes were right; the interiors were solid.",
           "size": 12.5, "color": MUTED, "spacing": 1.32, "after": 16},
          {"text": "The cause was structural. Each fragment picked its tone from its own "
                   "view's target, before anything was known about what else would occupy "
                   "the same beam — and with 93% of fragments serving all three views, "
                   "transmittances multiply. Mid-grey skin arrives as black.",
           "size": 12.5, "color": MUTED, "spacing": 1.32, "after": 16},
          {"text": "First fix — and it failed.", "size": 14, "color": FG, "after": 8},
          {"text": "Radiometric pre-compensation: divide the stray light out of the target "
                   "and re-solve. Standard practice in projector-camera systems. It "
                   "**oscillated** — 0.690 → 0.598 → 0.665 → 0.641 — because re-solving "
                   "changes which fragments exist, which changes the stray light, which "
                   "changes the target. A fixed point with no reason to contract.",
           "size": 12.5, "color": MUTED, "spacing": 1.32}])
    rect(s, 7.15, 1.9, 5.5, 2.95)
    text(s, 7.45, 2.1, 4.9, 2.55,
         [{"text": "What worked: freeze the geometry", "size": 14, "color": ACCENT,
           "after": 10},
          {"text": "Fragment shapes and host assignments are fixed. The only free variables "
                   "left are the **four discrete tone levels**. That turns an ill-posed fixed "
                   "point into plain **coordinate descent**: sweep one sheet at a time, try "
                   "all four tones, keep the best.",
           "size": 12, "color": MUTED, "spacing": 1.3, "after": 10},
          {"text": "Every sweep is monotone by construction, because a sweep only accepts a "
                   "change that lowers the loss.",
           "size": 12, "color": FG, "spacing": 1.3}])
    # Built from the shipped descent log, not from literals — the two columns exist
    # precisely to show that IoU barely registers the stage that fixed the faces.
    rows = [["Sweep", "composite score", "mean IoU", "pixels changed"]]
    for e in N["descent"]:
        rows.append([
            "start" if e["sweep"] < 0 else str(e["sweep"] + 1),
            f"{e['score']:.4f}",
            f"{e['mean_iou']:.4f}",
            "—" if e["sweep"] < 0 else f"{e['pixels_changed']:,}",
        ])
    grid_table(s, 7.15, 4.9, 5.5, rows, [1.0, 1.65, 1.35, 1.5], row_h=0.30)
    text(s, 7.15, 7.02, 5.5, 0.36,
         [{"text": f"score **+{N['score2'] - N['score1']:.4f}**  ·  mean IoU "
                   f"**+{N['iou2'] - N['iou1']:.4f}**  ·  the faces went from crushed to readable",
           "size": 10.5, "color": MUTED}])
    notes(s, f"""
[7:05-8:05] The second-strongest slide. It is the honest-failure slide, and panels reward it.

"Our renders hit {N['iou1']:.3f} on the standard metric and they looked wrong. Both of those were
true at the same time, and that is the interesting part. IoU scores the outline of a shape. A
face that is exactly the right shape and completely crushed to black scores well.

We tried the textbook fix first — divide the stray light out of the target and re-solve. It
oscillated. Re-solving changes which fragments exist, which changes the stray light, which
changes the target.

What worked was freezing the geometry entirely and only letting the tones move. That turns it
into coordinate descent, which is monotone by construction. Look at the pixel counts — each
sweep touches about a quarter of what the previous one did."

IF ASKED ABOUT THE TABLE — this is the sharpest point on the slide, so offer it:
the composite objective climbs {N['score2'] - N['score1']:+.4f} across the descent while mean IoU moves
{N['iou2'] - N['iou1']:+.4f}. The stage that rescued the faces is almost invisible to IoU. That is the
same blindness that hid the bug, showing up again while we fixed it — which is why we now
never report IoU without SSIM and an edge term beside it.
""")

    # ---- 11. fabrication -------------------------------------------------
    s = slide(prs, "Fabrication", "The project — fabrication")
    shots = [("fab_01_engraving", "Laser engraving the tone layers"),
             ("fab_02_sheets_flat", "The six sheets, cut and unassembled"),
             ("fab_03_assembly", "Slotting the weave together"),
             ("fab_04_object", "The finished assembly")]
    cw = (BODY_W - 3 * 0.24) / 4
    for i, (stem, cap) in enumerate(shots):
        photo(s, stem, MARGIN + i * (cw + 0.24), 1.95, cw, 3.05, caption=cap)
    rect(s, MARGIN, 5.12, 6.15, 1.78)
    text(s, MARGIN + 0.28, 5.36, 5.6, 1.35,
         [{"text": "ENG_L → ENG_D → ENG_K → CUT_SLOT → CUT_OUTLINE", "size": 11.5,
           "color": ACCENT, "after": 7},
          {"text": "Engrave first, cut last. Cutting the outline first leaves a compliant "
                   "part that lifts off the bed and defocuses the beam.",
           "size": 11, "color": MUTED, "spacing": 1.26, "after": 7},
          {"text": f"Widest slot **{N['slot']:.2f} mm** on 3 mm stock — not an error: a slot "
                   f"at a {N['angle']:.0f}° crossing must be thickness ÷ sin θ.",
           "size": 11, "color": MUTED, "spacing": 1.26}])
    rect(s, 7.42, 5.12, 5.19, 1.78)
    text(s, 7.7, 5.36, 4.65, 1.35,
         [{"text": "3 mm clear cast acrylic · 3,629 engraved regions · 12 joints",
           "size": 11.5, "color": ACCENT, "after": 7},
          {"text": "The four tones are engraving depths in **one** material — no dye, no "
                   "lamination, no thickness variation.",
           "size": 11, "color": MUTED, "spacing": 1.26, "after": 7},
          {"text": "Specified as **transmittance**, not as power and speed, because those "
                   "differ per machine. Calibrate on scrap first.",
           "size": 11, "color": MUTED, "spacing": 1.26}])
    notes(s, """
[8:05-8:55] Yahel's slide. Talk about the physical process, not the algorithm.

The two non-obvious things, and they are the ones fabrication people react to:

1. "Engrave before you cut. If you cut the outline first you are left with a floppy part that
   lifts off the bed and defocuses the beam."
2. "The tones are specified as transmittance, not as machine settings, so they have to be
   calibrated against a step wedge on scrap first. Correct files on an uncalibrated machine
   still give you a disappointing object."

The vectorisation pipeline — threshold, minimum feature, contour, kerf offset — is a BACKUP
slide. Go there only if asked how the raster becomes polygons.

TODO: drop four photos into examples/photos/ and re-run make_deck.py.
""")

    # ---- 12. object and shadows -----------------------------------------
    s = slide(prs, "The artefact, and the shadows it actually casts", "The project — result")
    photo(s, "fab_04_object", MARGIN, 1.95, 4.2, 4.55, caption="The assembly — roughly 30 × 30 cm")
    trio = [("shadow_front", "Front"), ("shadow_side", "Profile"), ("shadow_back", "Back")]
    cw = (BODY_W - 4.2 - 0.3 - 2 * 0.22) / 3
    for i, (stem, cap) in enumerate(trio):
        photo(s, stem, MARGIN + 4.5 + i * (cw + 0.22), 1.95, cw, 4.55, caption=cap)
    text(s, MARGIN, 6.65, BODY_W, 0.4,
         [{"text": "Same object, same lamp. Only the rotation changes.",
           "size": 13, "color": ACCENT}])
    notes(s, """
[8:55-9:30] Let this one sit. Say the caption line and stop talking for a beat.

If the room can be darkened, do this live with the object and a lamp instead of the slide —
it is far more convincing than any photograph. Check the room in advance.

TODO: drop the three shadow photos into examples/photos/.
""")

    # ---- 13. results ----------------------------------------------------
    s = slide(prs, "Results, and how we verified them", "The project — validation")
    cards(s, 1.9, [(f"{N['mean_iou']:.3f}", "mean IoU across the three views"),
                   (f"{N['min_iou']:.3f}", "worst single view — we optimise this too"),
                   (f"{N['ssim']:.3f}", "mean SSIM"),
                   ("−0.005", "cross-talk — constructive, not destructive")], height=1.4)
    text(s, MARGIN, 3.55, 6.1, 0.4,
         [{"text": "Three gates, all required to pass", "size": 14, "color": ACCENT}])
    gates = [("A — mechanical", "Every joint checked for crossing angle, slot width against "
                                "material thickness, and impossible triple junctions."),
             ("B — fabrication round trip", "Re-rendered from the exported polygons rather "
                                            "than the solver's raster. Cost: +0.0018 IoU."),
             ("C — ablation", "Each sheet deleted in turn. 0 of 6 idle; drops range "
                              "0.052 to 0.166.")]
    y = 4.0
    for head, body in gates:
        text(s, MARGIN, y, 6.1, 0.7,
             [{"text": head, "size": 12.5, "color": FG, "after": 4},
              {"text": body, "size": 11.5, "color": MUTED, "spacing": 1.28}])
        y += 0.82
    text(s, 7.2, 3.55, 5.45, 0.4,
         [{"text": "What the six-sheet constraint cost", "size": 14, "color": ACCENT}])
    grid_table(s, 7.2, 4.0, 5.45,
               [["", "6 sheets", "18 sheets"],
                ["mean IoU", f"{N['mean_iou']:.3f}", "0.882"],
                ["worst view", f"{N['min_iou']:.3f}", "0.871"],
                ["cut files", "13", "37"],
                ["joints to assemble", "12", "108"]],
               [2.55, 1.45, 1.45])
    text(s, 7.2, 5.95, 5.45, 0.8,
         [{"text": "Six sheets costs about 4% accuracy and removes **96 assembly joints**. "
                   "Every joint is an independent chance to misalign the whole piece, so we "
                   "would make that trade again.",
           "size": 11.5, "color": MUTED, "spacing": 1.3}])
    notes(s, """
[9:30-10:25]

Lead with the worst-view number, not the mean: "we optimise the worst of the three views
explicitly, because for a three-view piece the easiest way to cheat is to sacrifice one view
to flatter the average."

Gate C is worth thirty seconds — see the insights slide for the story about the sheet that
looked idle and was not.

The 6-vs-18 table pre-empts the obvious critique question ("why only six?"). Answer it before
they ask.
""")

    # ---- 14. it generalises ---------------------------------------------
    s = slide(prs, "The method is not specific to this painting", "Generality")
    gallery = [(ROOT / "out_final" / "faces2_60x60" / "preview_walls.png",
                "Three facial expressions — distinctness 0.97"),
               (ROOT / "out_final" / "cs_technion_60x60" / "preview_walls.png",
                "Two institutional logos"),
               (ROOT / "out_final" / "colour_wave" / "walls.png",
                "Hokusai series — subtractive colour mode")]
    cw = (BODY_W - 2 * 0.28) / 3
    for i, (path, cap) in enumerate(gallery):
        picture(s, path, MARGIN + i * (cw + 0.28), 2.0, cw, 3.5, caption=cap)
    text(s, MARGIN, 5.75, BODY_W, 1.0,
         [{"text": "The same solver produced all of these. Over the semester we ran it "
                   "across portraits, faces, logos, still lifes and landscapes, plus a "
                   "174-run study characterising cross-talk noise.",
           "size": 13, "color": FG, "after": 8},
          {"text": "The limits are known and documented: subjects must be centred on a plain "
                   "background, and full-frame paintings do not work — there is no silhouette "
                   "to carry.",
           "size": 11.5, "color": MUTED, "spacing": 1.3}])
    notes(s, """
[10:25-10:45] Twenty seconds. This is the first slide to cut if you are over time.

The honest caveat matters more than the gallery: "we know what it cannot do — the subject has
to be centred on a plain background. Full-frame paintings have no silhouette to carry, and
they fail."
""")

    # ---- 15. demo + video (template) -------------------------------------
    # Three required deliverables (parametric UI, CAD interop, the short film) collapsed onto
    # one slide, because at 15 minutes they are one minute between them.
    s = slide(prs, "Live demo, and the video", "Demo")
    cw = (BODY_W - 2 * 0.24) / 3
    xs_ = [MARGIN + i * (cw + 0.24) for i in range(3)]
    picture(s, _find("demo_ui"), xs_[0], 1.9, cw, 2.25,
            caption="Parametric solve, live or recorded", label="demo_ui.png")
    picture(s, _find("scene_3d"), xs_[1], 1.9, cw, 2.25,
            caption="The interactive 3D scene", label="scene_3d.png")
    picture(s, _find("video_still"), xs_[2], 1.9, cw, 2.25,
            caption="30–40 second summary video", label="video_still.png")
    panels = [
        ("Parametric interface",
         ["· sheet count and spacing",
          "· lamp distance and source radius → penumbra",
          "· fragment size and shard budget",
          "· the four engrave tone levels",
          "· stop angles and grid phase"],
         "Each is a field on the configuration object; a re-solve is 21.6 s."),
        ("CAD interoperability",
         ["· per-vertex-coloured **OBJ** and **PLY**",
          "· opens cleanly in Rhino, Blender, MeshLab",
          "· layered **DXF** / **SVG** per sheet, laser-ready",
          "· orbitable browser scene: sheets, lamp, rays,",
          "   and all three projected walls"],
         "No Grasshopper component yet — the solver is a Python library that a component "
         "would wrap."),
        ("The film",
         ["0:00  the six sheets, flat on the bench",
          "0:06  assembly — slotting the weave",
          "0:12  lamp on, the first portrait resolves",
          "0:20  slow rotation, morphing to the profile",
          "0:34  title card and names"],
         "The continuous rotation is the shot that sells the piece."),
    ]
    for x, (head, bullets, foot) in zip(xs_, panels):
        rect(s, x, 4.3, cw, 2.15)
        lines = [{"text": head, "size": 12.5, "color": ACCENT, "after": 7}]
        lines += [{"text": b, "size": 10, "color": FG, "spacing": 1.18, "after": 2}
                  for b in bullets]
        lines.append({"text": foot, "size": 9, "color": MUTED, "spacing": 1.2, "before": 5})
        text(s, x + 0.24, 4.48, cw - 0.48, 1.85, lines)
    template_banner(s)
    notes(s, """
[10:45-11:45] TEMPLATE — replace all three panels before presenting.

One minute for the whole slide. Do the demo live if you can; if the room or the laptop is
uncertain, play recordings. The fallback that always works is opening the interactive 3D
scene in a browser and orbiting it — it reads well on a projector.

Be direct about Grasshopper: there is no component. What exists is clean interop — OBJ, PLY,
DXF, SVG — plus our own browser viewer, which is exactly why we built it (slide 4). The panel
will respect "not yet, and here is what it would take" far more than a stretch.

The video plays from a SEPARATE FILE, not embedded, and it is also a required Aug 9
submission item.
""")

    # ---- 16. insights ---------------------------------------------------
    s = slide(prs, "What the semester actually taught us", "Insights")
    lessons = [
        ("A measured rule is only valid over the range it was measured on",
         "Both of our design rules reversed when we cut from eighteen sheets to six. Neither "
         "had been measured badly. Both had been measured while another variable was pinned "
         "by a constraint that then went away. This is a confounding failure, and it is not "
         "detectable from inside the original experiment."),
        ("A gate that only ever passes is not a gate",
         "Our ablation check earned its place the day it contradicted us. One sheet carried "
         "32 engraved regions against ~990 on its neighbours and looked like a blank. It has "
         "the **largest** ablation drop of all six. Region count is not contribution."),
        ("The right thing to differentiate is the wrong thing to accept on",
         "MSE is smooth and gives a usable descent direction, so we propose with it. "
         "Optimising it to convergence produces flat, low-contrast solutions that score well "
         "and look worse — so acceptance is judged on a composite objective instead."),
        ("Report a metric with the thing it is blind to",
         "IoU scores silhouettes. Across the entire re-toning stage — the stage that took the "
         f"faces from crushed to readable — it moved by **{N['iou2'] - N['iou1']:+.4f}**. It was "
         "not wrong; it was answering a different question. We now never report it without "
         "SSIM and an edge term beside it."),
    ]
    cw = (BODY_W - 0.26) / 2
    for i, (head, body) in enumerate(lessons):
        x = MARGIN + (i % 2) * (cw + 0.26)
        y = 1.92 + (i // 2) * 2.6
        rect(s, x, y, cw, 2.45)
        text(s, x + 0.3, y + 0.26, cw - 0.6, 1.95,
             [{"text": head, "size": 14, "color": ACCENT, "spacing": 1.15, "after": 10},
              {"text": body, "size": 11.5, "color": MUTED, "spacing": 1.3}])
    notes(s, """
[11:45-12:45] The closing argument. If you are over time, cut slides 5 and 14 to protect
this one — it is what the panel will remember.

Lead with the first. It pays off the Baran reversal you flagged on the prior-art slide:
"I said earlier that the rule we took from Baran reversed on us. That happened twice, with
two different rules, and the reason was the same both times."

The second one is the best anecdote in the project. Tell it as a story: "we looked at a sheet
with thirty-two regions next to sheets with nine hundred and assumed it was dead weight. It
turned out to be the single most important sheet in the piece."
""")

    # ---- 17. future ------------------------------------------------------
    s = slide(prs, "Where this goes next", "Future directions")
    future = [
        ("Look ahead, not just forward",
         "Host assignment is a single greedy pass — no lookahead, no revisiting. A fragment "
         "placed early cannot move once later ones reveal a better arrangement. Beam search, "
         "or a second reassignment pass conditioned on the stray field the first one actually "
         "produced, is the obvious next step."),
        ("Blue-noise tone quantisation",
         "Fragment boundaries still correlate weakly along the vertical axis, producing "
         "faint streaks in the result. Error diffusion with a blue-noise mask inside the "
         "re-toning stage is the standard halftoning answer, and is not yet implemented."),
        ("Continuous rotation",
         "Three stops is a design choice, not a limit. Supervising the intermediate angles "
         "would turn the piece from three portraits into a genuine animation — a face that "
         "turns, rather than three faces that swap."),
        ("A real parametric front end",
         "The configuration is already fully parametric and a solve takes 21.6 seconds. "
         "Sliders, and a Grasshopper component wrapping the same library, would put the "
         "design loop in a designer's hands rather than a programmer's — which is where this "
         "project started."),
    ]
    cw = (BODY_W - 3 * 0.24) / 4
    for i, (head, body) in enumerate(future):
        x = MARGIN + i * (cw + 0.24)
        rect(s, x, 1.95, cw, 3.7)
        text(s, x + 0.26, 2.2, cw - 0.52, 3.2,
             [{"text": f"0{i + 1}", "size": 11, "color": ACCENT, "after": 9},
              {"text": head, "size": 14.5, "color": ACCENT, "spacing": 1.14, "after": 11},
              {"text": body, "size": 11, "color": MUTED, "spacing": 1.3}])
    rect(s, MARGIN, 5.85, BODY_W, 1.05)
    text(s, MARGIN, 6.08, BODY_W, 0.7,
         [{"text": "Code, cut files, sweep journals and the full decision record:  "
                   "**github.com/shathan18/school**", "size": 13.5, "color": FG,
           "align": PP_ALIGN.CENTER, "after": 6},
          {"text": "Thank you — questions welcome.", "size": 11.5, "color": MUTED,
           "align": PP_ALIGN.CENTER}])
    notes(s, """
[12:45-13:15] Thirty seconds, then stop. Ten minutes of critique follows.

Close on the first item. Admitting the greedy is myopic is a strength, not a weakness, and it
is the question a panel is most likely to ask anyway — answering it before they do is worth
far more than the thirty seconds it costs.

The last item is the nicest closing line available: "we started wanting a designer's tool and
ended up writing a solver. Putting the solver back behind sliders is the last step."

THERE ARE BACKUP SLIDES after this one: the algorithm in full, the optics and rig numbers,
the vectorisation pipeline, the design sweep, and a list of the questions we expect. Know the
order so you can jump straight to one.
""")

    # ================================ backup ==============================
    # Everything below is off the 15-minute path. It exists so that a question can be
    # answered with evidence on screen instead of from memory.

    divider(prs, "Backup", "Not presented — held for the critique")
    notes(prs.slides[-1], """
Do not walk through these. Jump to one when a question lands:

  B1  the damage functional, its four limiting cases, and why greedy
  B2  the optics and the full rig geometry
  B3  raster to cut files — thresholding, minimum feature, kerf
  B4  the design sweep, and the two rules that reversed
  B5  the questions we expect, with the answers
""")

    # ---- B1. the algorithm in full ---------------------------------------
    s = slide(prs, "The damage functional, in full", "Backup — algorithm")
    text(s, MARGIN, 1.9, 6.55, 0.4,
         [{"text": "What a fragment costs the views it was not drawn for", "size": 13.5,
           "color": ACCENT}])
    rect(s, MARGIN, 2.4, 6.55, 1.15)
    text(s, MARGIN + 0.3, 2.6, 6.0, 0.85,
         [{"text": "\u03b4(f,p) = mean over q of   \u2016T_f \u2212 T_w(q)\u2016\u00b2  \u2212  "
                   "\u20161 \u2212 T_w(q)\u2016\u00b2", "size": 12, "color": FG,
           "font": MONO, "after": 7},
          {"text": "q = G(x),  for x in \u2264200 sampled pixels of fragment f",
           "size": 10.5, "color": MUTED, "font": MONO}])
    text(s, MARGIN, 3.68, 6.55, 0.75,
         [{"text": "T_f is the fragment's transmittance; T_w(q) is what view **w** wants at "
                   "the point its material lands on. The first term is the error **with** the "
                   "fragment there, the second the error of leaving that spot blank — so δ is "
                   "the change it causes, not its absolute darkness.",
           "size": 11, "color": MUTED, "spacing": 1.28}])
    cases = [("the other view wants…", "\u03b4", "consequence", True),
             ("white background", "> 0", "pure harm — a visible stray shard", False),
             ("darkness already", "< 0", "credit — the stray shadow does real work", False),
             ("\u2248 this fragment's tone", "< 0", "maximal credit", False),
             ("nothing — it lands off-wall", "0", "steering cross-talk away is learnable",
              False)]
    for i, (a, b, c, head) in enumerate(cases):
        yy = 4.52 + i * 0.36
        if head:
            rect(s, MARGIN, yy, 6.55, 0.36, PANEL, MSO_SHAPE.RECTANGLE)
        elif i % 2 == 0:
            rect(s, MARGIN, yy, 6.55, 0.36, RGBColor(0x17, 0x1A, 0x1F), MSO_SHAPE.RECTANGLE)
        col = MUTED if head else FG
        text(s, MARGIN + 0.16, yy + 0.07, 2.4, 0.24,
             [{"text": a, "size": 10.5 if head else 11.5, "color": col}])
        text(s, MARGIN + 2.6, yy + 0.07, 0.7, 0.24,
             [{"text": b, "size": 10.5 if head else 11.5,
               "color": MUTED if head else ACCENT, "font": None if head else MONO}])
        text(s, MARGIN + 3.45, yy + 0.07, 2.95, 0.24,
             [{"text": c, "size": 10.5 if head else 11, "color": MUTED}])
    text(s, MARGIN, 6.42, 6.55, 0.7,
         [{"text": "Because δ may be **negative**, the optimiser can be paid for cross-talk "
                   "rather than merely tolerating it. Measured on the shipped build: "
                   "**−0.005** — constructive.",
           "size": 11.5, "color": ACCENT, "spacing": 1.28}])
    rect(s, 7.42, 1.9, 5.19, 2.45)
    text(s, 7.7, 2.12, 4.65, 2.05,
         [{"text": "Why greedy, and what that costs", "size": 13, "color": ACCENT,
           "after": 9},
          {"text": "Assignment couples every pair of fragments through the multiplicative "
                   "composite, so the exact problem is **quadratic in the fragment count**. "
                   "We take one pass, ordered by demanded darkness, with no revisiting. It is "
                   "provably not optimal.",
           "size": 11, "color": MUTED, "spacing": 1.28, "after": 8},
          {"text": "The monotone second stage then recovers part of what the greedy left — "
                   "and that pair measured better than a worse solution to the exact problem "
                   "in the same wall-clock time.",
           "size": 11, "color": FG, "spacing": 1.28}])
    rect(s, 7.42, 4.5, 5.19, 1.55)
    text(s, 7.7, 4.72, 4.65, 1.15,
         [{"text": "Complexity", "size": 13, "color": ACCENT, "after": 8},
          {"text": "O(F · P · V · S),  S \u2264 200", "size": 11, "color": FG, "font": MONO,
           "after": 7},
          {"text": "F ≈ 1,168 fragments · P = 6 sheets · V = 3 views. Complete build 21.6 s "
                   "on an RTX 3060 Ti.",
           "size": 11, "color": MUTED, "spacing": 1.26}])
    text(s, 7.42, 6.2, 5.19, 0.8,
         [{"text": "And λ is not a dial.", "size": 12, "color": ACCENT, "after": 5},
          {"text": "Every non-zero damage weight gave **identical scores to four decimal "
                   "places**. It is a switch: price cross-talk, or do not.",
           "size": 11, "color": MUTED, "spacing": 1.26}])
    notes(s, """
BACKUP — go here for "how is damage actually defined?" or "is the greedy optimal?"

The functional is the thing to be proud of, because its four limiting cases fall out of the
algebra rather than being special-cased:

  - target white there   -> leaving it blank is already correct, so any shadow is pure harm
  - target dark there    -> blank is maximally WRONG, so a dark fragment IMPROVES it -> credit
  - target matches tone  -> the fragment supplies exactly the colour wanted -> maximal credit
  - lands off the wall   -> contributes nothing, so "aim the cross-talk into the void" is a
                            strategy the assignment can discover on its own

Earlier versions were harm-only and non-negative. That form structurally could not reward
double duty, so it drove good cross-talk to zero along with the bad. Making delta signed is
the single change that turned interference from a cost into a resource.
""")

    # ---- B2. optics and the rig ------------------------------------------
    s = slide(prs, "The optics fix what is physically knowable", "Backup — optics")
    text(s, MARGIN, 1.95, 5.8, 4.3,
         [{"text": "Central projection → one 3×3 matrix", "size": 15, "color": ACCENT,
           "after": 7},
          {"text": "A point light, a flat sheet and a flat wall define a projective "
                   "collineation, so the whole mapping is a **homography**. We build it from "
                   "the four sheet corners by normalised DLT and SVD, and cache it per "
                   "(sheet, view) pair.",
           "size": 12.5, "color": MUTED, "spacing": 1.32, "after": 20},
          {"text": "Magnification is the depth-to-resolution gradient", "size": 15,
           "color": ACCENT, "after": 7},
          {"text": "A 29.6 cm sheet becomes a **1.80 m** image — **6.09×**. Sheets nearer the "
                   "lamp throw larger, coarser shadows; sheets nearer the wall carry finer "
                   "detail. Depth is therefore also a detail budget.",
           "size": 12.5, "color": MUTED, "spacing": 1.32, "after": 20},
          {"text": "Penumbra is a hard resolution bound", "size": 15, "color": ACCENT,
           "after": 7},
          {"text": "A real lamp has width, so every shadow edge is blurred — **3.6 to "
                   "3.9 mm** at the wall here. We compute that first and refuse to engrave "
                   "anything finer. Detail below the blur is precision the light can never "
                   "resolve; fabricating it is a lie told in acrylic.",
           "size": 12.5, "color": MUTED, "spacing": 1.32}])
    rect(s, 6.9, 1.9, 5.72, 4.55)
    text(s, 7.2, 2.1, 5.15, 0.4,
         [{"text": "The rig these numbers describe", "size": 13.5, "color": ACCENT}])
    geo = [("lamp → body", "0.493 m"),
           ("body → wall", "2.507 m"),
           ("lamp height", "0.874 m"),
           ("sheet", "295.8 × 295.8 mm, 3 mm cast acrylic"),
           ("footprint", "298.7 × 298.7 mm"),
           ("image on the wall", "1.80 m"),
           ("magnification", "6.085×"),
           ("penumbra at the wall", "3.59 / 3.93 / 3.93 mm"),
           ("smallest honest feature", "30.4 mm at the wall"),
           ("turntable stops", "15° / 135° / 255°")]
    for i, (k, v) in enumerate(geo):
        y = 2.62 + i * 0.375
        text(s, 7.2, y, 2.5, 0.34, [{"text": k, "size": 11.5, "color": MUTED}])
        text(s, 9.7, y, 2.62, 0.34,
             [{"text": v, "size": 11.5, "color": FG, "align": PP_ALIGN.RIGHT}])
    notes(s, """
BACKUP — go here for anything about the rig, the blur, or "what resolution can you actually
achieve?"

The line worth having ready: "we compute the blur from the lamp's physical size first, and
then refuse to engrave anything finer than it. Detail below the blur is precision the light
can never resolve."

Those ten numbers plus the cut files are sufficient to rebuild the piece, which is why they
are also in the replication section of the submission document.
""")

    # ---- B3. raster to cut -----------------------------------------------
    s = slide(prs, "From a solved field to files a laser can run", "Backup — fabrication")
    steps = [
        ("Threshold", "Each sheet's solved field is split into the four tone levels."),
        ("Minimum feature", "Morphological opening then closing at the cuttable radius — "
                            "this is where the penumbra bound becomes physical geometry."),
        ("Contour", "Regions become polygons in sheet-local millimetres."),
        ("Kerf offset", "Every ring is grown by half the beam width, so the finished part is "
                        "nominal size after the laser eats the edge."),
    ]
    text(s, MARGIN, 1.95, 5.7, 0.4,
         [{"text": "Vectorisation", "size": 15, "color": ACCENT}])
    y = 2.45
    for i, (head, body) in enumerate(steps):
        text(s, MARGIN, y, 5.7, 0.8,
             [{"text": f"{i + 1}.  {head}", "size": 13, "color": FG, "after": 4},
              {"text": body, "size": 11.5, "color": MUTED, "spacing": 1.28}])
        y += 0.86
    rect(s, MARGIN, 5.95, 5.7, 0.95)
    text(s, MARGIN + 0.28, 6.18, 5.2, 0.6,
         [{"text": "Layer order: ENG_L → ENG_D → ENG_K → CUT_SLOT → CUT_OUTLINE",
           "size": 12, "color": ACCENT, "after": 5},
          {"text": "Engrave first, cut last. Cutting the outline first leaves a compliant "
                   "part that lifts and defocuses.", "size": 11, "color": MUTED}])
    picture(s, ROOT / "examples" / "pearl3" / "4_fabrication.png", 6.85, 1.9, 5.77, 3.4)
    facts = [["6", "sheets, 295.8 mm square, 3 mm clear acrylic"],
             ["3,629", "engraved regions"],
             ["4", "tone levels — clear, 0.85, 0.62, 0.35 transmittance"],
             ["12", "half-lap joints"],
             [f"{N['slot']:.2f} mm", f"widest slot on 3.0 mm stock, at {N['angle']:.0f}° crossing"]]
    yy = 5.45
    for value, label in facts:
        text(s, 6.85, yy, 1.5, 0.28, [{"text": value, "size": 12.5, "color": ACCENT}])
        text(s, 8.35, yy, 4.27, 0.28, [{"text": label, "size": 11.5, "color": MUTED}])
        yy += 0.29
    notes(s, """
BACKUP — go here for "how does the solved image become polygons?" or any kerf / tolerance
question.

The minimum-feature step is the interesting one: it is where the penumbra bound computed on
the optics slide stops being a number and becomes physical geometry. We open then close at
the cuttable radius, so nothing survives that the light could not have resolved anyway.
""")

    # ---- B4. the design sweep --------------------------------------------
    s = slide(prs, "The design sweep, and the two rules that reversed", "Backup — evidence")
    picture(s, ROOT / "examples" / "pearl3" / "2_levers.png", MARGIN, 1.85, 7.55, 5.0)
    text(s, 8.5, 1.9, 4.11, 5.0,
         [{"text": "We swept the design, not the solver.", "size": 14.5, "color": ACCENT,
           "after": 9},
          {"text": "Sheet count, pitch, engrave alphabet, damage weight and fragment size — "
                   "**two seeds per point**, accepted on a composite of mean IoU, worst view "
                   "and SSIM rather than on IoU alone.",
           "size": 11.5, "color": MUTED, "spacing": 1.3, "after": 18},
          {"text": "Pitch wanted to open, not tighten.", "size": 13, "color": FG, "after": 7},
          {"text": "Tight stacking was right when sheets were many. At six, **50 mm** takes "
                   "the worst view from 0.785 to 0.816 and drives cross-talk to zero. "
                   "100 mm collapses it again.",
           "size": 11.5, "color": MUTED, "spacing": 1.3, "after": 18},
          {"text": "The engrave alphabet wanted to go light.", "size": 13, "color": FG,
           "after": 7},
          {"text": "Eighteen pale layers multiply into darkness, so the tones must be dark. "
                   "With two layers per view there is little left to multiply, and a dark "
                   "alphabet overshoots. We ship **0.85 / 0.62 / 0.35** — the rule from the "
                   "literature finished **sixth of seven**.",
           "size": 11.5, "color": MUTED, "spacing": 1.3}])
    notes(s, """
BACKUP — this is the evidence for the reversal claimed on the prior-art slide and paid off on
the insights slide. Go here if anyone challenges either.

Point at the bottom-left panel: the green bar is what we shipped; the dark, density-even
alphabet from Baran is the fourth bar down.

The bottom-right panel is the whole semester in one chart — a third of the sheets and a
quarter of the footprint of our earlier build, and still ahead of it. The last bar is measured
from the exported cut files, not from the solver.

The honest framing: neither rule was measured badly. Both were measured at eighteen sheets,
where a constraint we later removed was doing the work. You cannot detect that from inside
the original experiment.
""")

    # ---- B5. expected questions ------------------------------------------
    s = slide(prs, "Questions we expect", "Backup — Q&A")
    qa = [
        ("Why only six sheets?",
         "Eighteen scores 0.882 against our 0.846 — about 4%. It also costs 24 more cut "
         "files and **96 more assembly joints**, and every joint is an independent chance to "
         "misalign the whole piece."),
        ("Is the greedy optimal?",
         "No. One pass, no revisiting; the exact problem is quadratic in the fragment count. "
         "Beam search or a second reassignment pass is the first item of future work."),
        ("How do you know the files match the simulation?",
         "Gate B re-renders from the **exported polygons**, not the solver's own raster. The "
         "round trip costs +0.0018 mean IoU, and that number is measured every build."),
        ("What if a sheet is mis-cut?",
         "The piece fails rather than degrades. At six sheets there is no redundancy — "
         "ablation drops range 0.052 to 0.166 — and that is a deliberate trade."),
        ("Why not three separate stacks?",
         "Then it is three objects in a box, not one object. That every piece of matter is "
         "visible to all three views at once **is** the problem; removing it removes the "
         "project."),
        ("Does it work for any three images?",
         "No. Subjects must be centred on a plain background; full-frame paintings have no "
         "silhouette to carry and fail. The distinctness of the three targets predicts it."),
    ]
    cw = (BODY_W - 2 * 0.24) / 3
    for i, (q, a) in enumerate(qa):
        x = MARGIN + (i % 3) * (cw + 0.24)
        y = 1.95 + (i // 3) * 2.45
        rect(s, x, y, cw, 2.3)
        text(s, x + 0.26, y + 0.24, cw - 0.52, 1.85,
             [{"text": q, "size": 13, "color": ACCENT, "spacing": 1.14, "after": 10},
              {"text": a, "size": 11, "color": MUTED, "spacing": 1.3}])
    notes(s, """
BACKUP — keep this slide up during the critique if the panel is asking rapid-fire questions.

Whoever is not answering should be finding the right backup slide.
""")

    prs.save(OUT)
    return prs


if __name__ == "__main__":
    prs = build()
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    if _missing:
        print("\nplaceholders still to fill — drop files into examples/photos/:")
        for name in dict.fromkeys(_missing):
            print(f"  {name}")
    else:
        print("\nall images resolved.")
